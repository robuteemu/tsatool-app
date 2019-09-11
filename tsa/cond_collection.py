#!/usr/bin/python
# -*- coding: utf-8 -*-

# Collection of Conditions for analysis

# TODO: handling Postgres connection instances
# is very ambiguous at the moment. Either the connection
# should be treated robustly as an attribute or
# given always from outside as an argument.

import logging
import traceback
import pptx
import openpyxl as xl
from .condition import Condition
from .utils import strfdelta
from .utils import list_local_statids
from .utils import list_local_sensors
from .tsaerror import TsaError
from datetime import datetime
from io import BytesIO
from collections import OrderedDict
from pptx.util import Pt
from pptx.util import Cm
from pptx.dml.color import RGBColor

log = logging.getLogger(__name__)

class CondCollection:
    """
    A collection of conditions to analyze.
    All conditions share the same analysis time range.
    Times are assumed to be given as dates only,
    and their HH:MM:SS are overridden to default values,
    see ``set_default_times(self)``.

    :param time_from: start time (inclusive) of the analysis period
    :type time_from: Python ``datetime()`` object
    :param time_until: end time (exclusive) of the analysis period
    :type time_until: Python ``datetime()`` object
    :param pg_conn: database connection
    :type pg_conn: ``psycopg2.connect()`` object
    """
    def __init__(self, time_from, time_until, pg_conn=None, title=None):
        # Times must be datetime objects and in correct order
        assert isinstance(time_from, datetime)
        assert isinstance(time_until, datetime)
        assert time_from <= time_until
        self.time_from = time_from
        self.time_until = time_until
        self.set_default_times()
        self.time_range = (self.time_from, self.time_until)

        self.title = title

        # Timestamp is based on instance creation time,
        # not on when the analysis has been run
        self.created_timestamp = datetime.now()

        self.conditions = OrderedDict()
        self.station_ids = set()

        self.errors = []

        self.pg_conn = pg_conn
        self.statids_available = set()
        self.temptables = []

    def setup_views(self):
        """
        Set up time-limited statobs view and joint main observation view.
        """
        self.setup_statobs_view()
        self.setup_obs_view()

    def add_error(self, msg, lvl='Error'):
        """
        Add error message to error message list.
        Only unique errors are collected, in order to avoid
        piling up repetitive messages from loops, for example.
        """
        err = TsaError(lvl=lvl, cxt=f'Sheet {self.name}', msg=msg)
        if err not in self.errors:
            self.errors.append(err)

    def set_default_times(self):
        """
        Sets analysis start time to 00:00:00
        and end time to 23:59:59 on selected dates, respectively.
        """
        self.time_from = self.time_from.replace(
            hour=0, minute=0, second=0)
        self.time_until = self.time_until.replace(
            hour=23, minute=59, second=59)

    def setup_statobs_view(self, verbose=False):
        """
        In the database, create or replace a temporary view ``statobs_time``
        containing the station observations within the ``time_range``.
        """
        if self.pg_conn:
            try:
                with self.pg_conn.cursor() as cur:
                    sql = ("CREATE OR REPLACE TEMP VIEW statobs_time AS "
                           "SELECT id, tfrom, statid "
                           "FROM statobs "
                           "WHERE tfrom BETWEEN %s AND %s;")
                    if verbose:
                        log.debug(cur.mogrify(sql, (self.time_from, self.time_until)))
                    cur.execute(sql, (self.time_from, self.time_until))
                    self.pg_conn.commit()
            except:
                self.pg_conn.rollback()
                msg = 'Could not create db view "statobs_time"'
                log.error(msg, exc_info=True)
                self.add_error(msg=msg)
        else:
            msg = 'No db connection, cannot create view "statobs_time"'
            log.warning(msg)
            self.add_error(msg)

    def get_stations_in_view(self):
        """
        Get stations available in ``statobs_time`` view.
        """
        if self.pg_conn:
            try:
                with self.pg_conn.cursor() as cur:
                    sql = "SELECT DISTINCT statid FROM statobs_time ORDER BY statid;"
                    cur.execute(sql)
                    statids = cur.fetchall()
                    statids = [el[0] for el in statids]
                    self.statids_available = set(statids)
            except:
                self.pg_conn.rollback()
                msg = 'Could not get station ids from db view "statobs_time"'
                log.error(f'Collection {self.title}: {msg}', exc_info=True)
                self.add_error(msg=msg)
        else:
            msg = 'No db connection, cannot get station ids from db view "statobs_time'
            log.error(f'Collection {self.title}: {msg}')
            self.add_error(msg)

    def setup_obs_view(self, verbose=False):
        """
        After creating the ``statobs_time`` view,
        create a joint temporary view ``obs_main``
        that works as the main source for Block queries.
        """
        if not self.pg_conn:
            msg = 'No db connection, cannot set up view "obs_main"'
            log.error(f'Collection {self.title}: {msg}')
            self.add_error(msg)
            return
        try:
            with self.pg_conn.cursor() as cur:
                sql = ("CREATE OR REPLACE TEMP VIEW obs_main AS "
                       "SELECT tfrom, statid, seid, seval "
                       "FROM statobs_time "
                       "INNER JOIN seobs "
                       "ON statobs_time.id = seobs.obsid;")
                if verbose:
                    log.debug(sql)
                cur.execute(sql)
                self.pg_conn.commit()
        except:
            self.pg_conn.rollback()
            msg = 'Could not create db view "obs_main"'
            log.error(f'Collection {self.title}: {msg}', exc_info=True)
            self.add_error(msg=msg)

    def add_condition(self, site, master_alias, raw_condition, excel_row=None):
        """
        Add new Condition instance, add error and skip if one exists already
        with same site-master_alias identifier.
        """
        try:
            candidate = Condition(site, master_alias, raw_condition, self.time_range, excel_row)
            if candidate.id_string in self.condtitions.keys():
                msg = (f'Site-master_alias combo {candidate.id_string} already reserved, '
                       'cannot add it twice')
                if excel_row is not None:
                    msg += f' (row {excel_row} in Excel)'
                log.error(f'Collection {self.title}: {msg}', exc_info=True)
                self.add_error(msg)
            else:
                for stid in candidate.station_ids:
                    self.station_ids.add(stid)
                    if stid not in self.statids_available:
                        msg = f'Station id {stid} of condition {candidate.id_string} not in available station ids'
                        log.warning(msg)
                        candidate.add_error(msg=msg, lvl='Warning')
                self.conditions[candidate.id_string] = candidate
        except:
            msg = 'Could not add condition'
            if excel_row is not None:
                msg += f' (row {excel_row} in Excel)'
            log.error(f'Collection {self.title}: {msg}', exc_info=True)
            self.add_error(msg)

    def set_sensor_ids(self, pairs=None):
        """
        Get sensor name - id pairs from the database,
        and set sensor ids for all Blocks in all Conditions.
        Optionally, the ``nameids`` can be fed from outside, in which case
        querying the database is omitted.
        """
        if pairs is None or len(pairs) == 0:
            if not self.pg_conn:
                self.add_error('WARNING: No db connection, cannot get sensor ids from database')
                return
            with self.pg_conn.cursor() as cur:
                cur.execute("SELECT id, lower(name) AS name FROM sensors;")
                tb = cur.fetchall()
                pairs = {k:v for v, k in tb}
        for cndk in self.conditions.keys():
            for i in range(len(self.conditions[cndk].blocks)):
                try:
                    self.conditions[cndk].blocks[i].set_sensor_id(pairs)
                except:
                    bl = self.conditions[cndk].blocks[i].alias
                    msg = f'Could not set sensor name-id pairs for block {f}'
                    log.error(msg, exc_info=True)
                    self.conditions[cndk].add_error(msg)

    def get_temporary_relations(self):
        """
        Set str list of temporary tables and views currently available in db.
        """
        if not self.pg_conn:
            self.add_error('WARNING: No db connection, cannot get temporary relations list')
            return
        with self.pg_conn.cursor() as cur:
            try:
                sql = ("SELECT table_name FROM information_schema.tables "
                       "WHERE table_schema LIKE '%pg_temp%';")
                cur.execute(sql)
                res = cur.fetchall()
            except Exception as e:
                self.pg_conn.rollback()
                self.add_error(e)
                return
        self.temptables = [el[0] for el in res]

    def create_condition_temptables(self, verbose=False):
        """
        For each Condition, create the corresponding temporary table in db.
        Primary conditions are handled first, only then secondary ones;
        if there are secondary conditions depending further on each other,
        it is up to the user to give them in correct order!
        """

        # First round for primary ones only
        for cndk in self.conditions.keys():
            if self.conditions[cndk].secondary:
                continue
            try:
                cnd.create_db_temptable(pg_conn=self.pg_conn,
                                        verbose=verbose,
                                        src_tables=self.temptables)
            except Exception as e:
                log.exception(e)
        self.get_temporary_relations()

        # Second round for secondary ones,
        # viewnames list is now updated every time
        for cndk in self.conditions.keys():
            if self.conditions[cndk].secondary:
                self.conditions[cndk].create_db_temptable(
                    pg_conn=self.pg_conn,
                    verbose=verbose,
                    src_tables=self.temptables)
                self.get_temporary_relations()

    def fetch_all_results(self):
        """
        Fetch results
        for all Conditions that have a corresponding view in the database.
        """
        cnd_len = len(self.conditions)
        for i, cndk in enumerate(self.conditions.keys()):
            log.debug(f'Fetching {i+1}/{cnd_len} {self.conditions[cndk].id_string} ...')
            self.conditions[cndk].fetch_results_from_db(pg_conn=self.pg_conn)

    def to_worksheet(self, wb):
        """
        Add a worksheet to an ``openpyxl.Workbook`` instance
        containing summary results of the condition collection.
        """
        assert isinstance(wb, xl.Workbook)
        ws = wb.create_sheet()
        ws.title = self.title or 'conditions'

        # Headers in fixed cells & styling
        headers = {'A1': 'start',
                   'B1': 'end',
                   'D1': 'analyzed',
                   'A3': 'site',
                   'B3': 'master_alias',
                   'C3': 'condition',
                   'D3': 'data_from',
                   'E3': 'data_until',
                   'F3': 'valid',
                   'G3': 'notvalid',
                   'H3': 'nodata',
                   'I3': 'rows'
                   }
        for k, v in headers.items():
            ws[k] = v
            ws[k].font = xl.styles.Font(bold=True)

        # Global values
        ws['A2'] = self.time_from
        ws['B2'] = self.time_until
        ws['D2'] = self.created_timestamp

        # Condition rows
        r = 4
        for cnd in self.conditions.values():
            ws[f'A{r}'] = cnd.site
            ws[f'B{r}'] = cnd.master_alias
            ws[f'C{r}'] = cnd.condition
            ws[f'D{r}'] = cnd.data_from
            ws[f'E{r}'] = cnd.data_until
            ws[f'F{r}'] = cnd.percentage_valid
            ws[f'G{r}'] = cnd.percentage_notvalid
            ws[f'H{r}'] = cnd.percentage_nodata
            ws[f'I{r}'] = cnd.n_rows

            # Percent format
            ws[f'F{r}'].number_format = '0.00 %'
            ws[f'G{r}'].number_format = '0.00 %'
            ws[f'H{r}'].number_format = '0.00 %'

            r += 1

    def to_pptx(self, pptx_template):
        """
        Return a ``pptx`` presentation object,
        making a slide of each condition.

        ``pptx`` must be a filepath or file-like object
        representing a PowerPoint file that includes the master
        layout for the TSA report and nothing else. The default
        placeholder indices must conform with the constants here!
        """
        phi = dict(
        HEADER_IDX = 17,     # Slide header placeholder
        TITLE_IDX = 0,       # Condition title placeholder
        BODY_IDX = 13,       # Condition string placeholder
        TIMERANGE_IDX = 15,  # Placeholder for condition start/end time text
        VALIDTABLE_IDX = 18, # Validity time/percentage table placeholder
        ERRORS_IDX = 19,     # Placeholder for errors and warnings
        MAINPLOT_IDX = 11,   # Main timeline plot placeholder
        FOOTER_IDX = 16,     # Slide footer placeholder
        )
        MAINPLOT_H_PX = 3840 # Main timeline plot height in pixels

        pres = pptx.Presentation(pptx_template)
        layout = pres.slide_layouts[0]

        # Ensure placeholder indices exist as they should
        indices_in_pres = [ph.placeholder_format.idx for ph in layout.placeholders]
        for k, v in phi.items():
            if v not in indices_in_pres:
                raise Exception(f'{k} {v} not in default layout placeholders')

        # Add slides and fill in contents for each condition.
        for c in self.conditions.values():
            s = pres.slides.add_slide(layout)

            # Slide header
            txt = 'TSA report: '
            if self.title is not None:
                txt += self.title
            txt += ' ' + self.created_timestamp.strftime('%d.%m.%Y')
            s.placeholders[phi['HEADER_IDX']].text = txt

            # Slide footer
            txt = 'TSATool v0.1, copyright WSP Finland'
            s.placeholders[phi['FOOTER_IDX']].text = txt

            # Condition title
            s.placeholders[phi['TITLE_IDX']].text = c.id_string

            # Condition string / body
            s.placeholders[phi['BODY_IDX']].text = c.condition

            # Condition data time range
            if not (c.data_from is None or c.data_until is None):
                txt = 'Datan tarkasteluväli {}-{}'.format(
                    c.data_from.strftime('%d.%m.%Y %H:%M'),
                    c.data_until.strftime('%d.%m.%Y %H:%M')
                )
            else:
                txt = 'Ei dataa saatavilla'
            s.placeholders[phi['TIMERANGE_IDX']].text = txt

            # Master condition validity table
            tb_shape = s.placeholders[phi['VALIDTABLE_IDX']].insert_table(rows=3, cols=4)
            tb = tb_shape.table

            tb.cell(0, 0).text = ''

            tb.cell(0, 1).text = 'Voimassa'
            tb.cell(0, 2).text = 'Ei voimassa'
            tb.cell(0, 3).text = 'Tieto puuttuu'

            tb.cell(1, 0).text = 'Yhteensä'
            txt = strfdelta(c.tottime_valid, '{days} pv {hours} h {minutes} min')
            tb.cell(1, 1).text = txt
            txt = strfdelta(c.tottime_notvalid, '{days} pv {hours} h {minutes} min')
            tb.cell(1, 2).text = txt
            txt = strfdelta(c.tottime_nodata, '{days} pv {hours} h {minutes} min')
            tb.cell(1, 3).text = txt

            tb.cell(2, 0).text = 'Osuus tarkasteluajasta'
            txt = '{} %'.format(round(c.percentage_valid*100, 2))
            tb.cell(2, 1).text = txt
            txt = '{} %'.format(round(c.percentage_notvalid*100, 2))
            tb.cell(2, 2).text = txt
            txt = '{} %'.format(round(c.percentage_nodata*100, 2))
            tb.cell(2, 3).text = txt

            for cl in tb.iter_cells():
                cl.fill.background()
                for ph in cl.text_frame.paragraphs:
                    ph.font.name = 'Montserrat'
                    ph.font.size = Pt(8)
                    ph.font.color.rgb = RGBColor.from_string('000000')

            for row in tb.rows:
                row.height = Cm(0.64)

            # Condition errors and warnings
            txt = '; '.join([str(msg) for msg in c.errmsgs]) or ' '
            s.placeholders[phi['ERRORS_IDX']].text = txt

            # Condition main timeline plot; ignored if no data to viz
            if c.main_df is None:
                continue
            # Find out the proportion of plot height of the width
            wh_factor = s.placeholders[phi['MAINPLOT_IDX']].height \
                        / s.placeholders[phi['MAINPLOT_IDX']].width
            w, h = MAINPLOT_H_PX, wh_factor*MAINPLOT_H_PX
            with BytesIO() as fobj:
                c.save_timelineplot(fobj, w, h)
                s.placeholders[phi['MAINPLOT_IDX']].insert_picture(fobj)

        return pres

    def save_pptx(self, pptx_template, out_path):
        """
        Call ``.to_pptx`` and save result to file.
        """
        pptx_obj = self.to_pptx(pptx_template=pptx_template)
        pptx_obj.save(out_path)

    def run_analysis(self, pg_conn, wb=None, pptx_path=None, pptx_template=None):
        """
        Call necessary methods to run the condition analysis
        and save results to the specified
        ``openpyxl.Workbook`` instance ``wb`` as new worksheet
        and the ``pptx_path`` as ``.pptx`` file.
        If an output is ``None``, it is not created.
        """
        log.info(f'Started analysis for collection {self.title}')
        self.pg_conn = pg_conn
        log.info('Setting up DB views')
        self.setup_views()
        log.info('Creating condition views')
        self.create_condition_temptables()
        self.fetch_all_results()

        if wb is not None:
            try:
                log.info(f'Adding {self.title} to Excel workbook')
                self.to_worksheet(wb)
            except:
                log.exception('Could not make Excel report sheet')

        if pptx_path is not None and pptx_template is not None:
            try:
                log.info(f'Saving pptx report to {pptx_path}')
                self.save_pptx(pptx_template=pptx_template, out_path=pptx_path)
            except:
                log.exception('Could not make pptx report')
        log.info(f'END OF ANALYSIS for collection {self.title}')

    def __getitem__(self, key):
        """
        Returns the Condition instance on the corresponding index.
        """
        return self.conditions[key]

    def __str__(self):
        t = self.title or '(no title)'
        s = f'Collection {t} with {len(self.conditions) conditions}'
        return s

    @classmethod
    def from_xlsx_sheet(cls, ws,
                        pg_conn=None,
                        station_ids=None,
                        sensor_pairs=None):
        """
        Create a condition collection for analysis
        based on an ``openpyxl`` ``worksheet`` object ``ws``.
        Database connection instance ``pg_conn`` should be prepared
        in advance and passed to this method.

        .. note:: Start and end dates must be in cells A2 and B2, respectively,
                  and conditions must be listed starting from row 4,
                  such that ``site`` is in column A,
                  ``master_alias`` in column B
                  and ``raw_condition`` in column C.
                  There must not be empty rows in between.
                  Any columns outside A:C are ignored,
                  so additional data can be placed outside them.
        """
        # Validate start and end dates.
        # These must be d.m.Y dates, start in cell A2
        # and end in cell B2. On error, record error but
        # use today as default date so constructing the
        # collection can go further.
        dateformat = '%d.%m.%Y'
        time_errs = [] # Supply this to the instance after it is created
        time_from = ws['A2'].value
        if time_from is None:
            time_from = datetime.now()
            time_errs.append(("Start date in cell A2 is empty: "
                              "must be a d.m.YYYY date"))
        if not isinstance(time_from, datetime):
            try:
                time_from = datetime.strptime(ws['A2'].value, dateformat)
            except:
                time_from = datetime.now()
                time_errs.append(("Could not read start date in cell A2: "
                                  "must be a d.m.YYYY date"))
        time_until = ws['B2'].value
        if time_until is None:
            time_until = datetime.now()
            time_errs.append(("End date in cell B2 is empty: "
                              "must be a d.m.YYYY date"))
        if not isinstance(time_until, datetime):
            try:
                time_until = datetime.strptime(ws['B2'].value, dateformat)
            except:
                time_from = datetime.now()
                time_errs.append(("Could not read end date in cell B2: ",
                                  "must be a d.m.YYYY date"))
        if time_from > time_until:
            time_errs.append("Start date (A2) must be BEFORE end date (B2)")
            # Set time_from so that further validation is possible
            time_from = time_until

        # With the start & end times prepared, initialize the instance
        # and request available station ids / set them from argument.
        # Then add conditions row by row, possibly adding any errors related.
        # Having made the conditions, set the sensor ids from database
        # or argument provided here; error is raised for each Block
        # if no sensor id is found.
        cc = cls(time_from=time_from, time_until=time_until,
                 pg_conn=pg_conn, title=ws.title)
        for terr in time_errs:
            cc.add_error(terr)
        if station_ids is not None:
            cc.statids_available = set(station_ids)
        empty_cells = []
        for row in ws.iter_rows(min_row=4, max_col=3):
            cells = [c for c in row]
            cells_ok = True
            for c in cells:
                if c.value is None:
                    empty_cells.append(c)
                    cells_ok = False
            if not cells_ok:
                continue
            cc.add_condition(site=cells[0].value, master_alias=cells[1].value,
                             raw_condition=cells[2].value, excel_row=cells[0].row)
        last_row = cells[0].row
        empty_cells = [c for c in empty_cells if c.row < last_row]
        for ec in empty_cells:
            cc.add_error(f"Cell {c.coordinate} should not be empty: row ignored")
        if len(cc.statids_available) == 0:
            cc.get_stations_in_view() # This requires valid pg_conn
        # If pairs is not given above as argument,
        # this will require valid pg_conn:
        cc.set_sensor_ids(pairs=sensor_pairs)

        return cc
